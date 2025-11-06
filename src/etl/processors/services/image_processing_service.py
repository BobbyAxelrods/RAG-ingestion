"""
Image Processing Service.

Handles extraction and enrichment of images from documents AFTER
Document Intelligence processing. This includes:
- Extracting embedded images from PDFs/Office docs
- Applying OCR to images
- Generating descriptions using GPT-4 Vision
- Enriching markdown with image context
"""

import logging
import re
from io import BytesIO

from PIL import Image

from src.config import ImageProcessingConfig
from src.models.document_models import EnrichedDocument, ImageInfo, PageContent, DocumentType
from src.services.doc_intel_service import DocumentIntelligenceService
from src.services.openai_service import OpenAIService

logger = logging.getLogger(__name__)


class ImageProcessingService:
    """
    Service for processing images from documents.

    Extracts images identified by Document Intelligence, applies OCR,
    generates descriptions, and enriches document content.
    """

    def __init__(
        self,
        config: ImageProcessingConfig,
        doc_intel_service: DocumentIntelligenceService,
        openai_service: OpenAIService,
    ):
        """
        Initialize Image Processing service.

        Args:
            config: Image processing configuration
            doc_intel_service: Document Intelligence service for OCR
            openai_service: OpenAI service for image descriptions
        """
        self.config = config
        self.doc_intel_service = doc_intel_service
        self.openai_service = openai_service

        logger.info(
            f"Initialized ImageProcessingService "
            f"(extract_embedded={config.extract_embedded_images}, "
            f"ocr={config.ocr_enabled}, summary_method={config.summary_method})"
        )

    def process_document_images(
        self, enriched_doc: EnrichedDocument, original_bytes: bytes
    ) -> EnrichedDocument:
        """
        Process all images in a document.

        This is the main orchestration method that:
        1. Extracts embedded images
        2. Applies OCR to each image
        3. Generates descriptions
        4. Enriches markdown with image context

        Args:
            enriched_doc: Document after Document Intelligence processing
            original_bytes: Original document bytes (for image extraction)

        Returns:
            EnrichedDocument: Document with processed images and enriched markdown

        Example:
            >>> image_service = ImageProcessingService(config, doc_intel, openai)
            >>> enriched_doc = image_service.process_document_images(doc, pdf_bytes)
            >>> print(f"Processed {len(enriched_doc.images)} images")
        """
        if not self.config.extract_embedded_images or not enriched_doc.images:
            logger.info(
                f"No images to process for {enriched_doc.filename} "
                f"(extract={self.config.extract_embedded_images}, "
                f"images={len(enriched_doc.images)})"
            )
            return enriched_doc

        logger.info(
            f"Processing {len(enriched_doc.images)} images from {enriched_doc.filename}"
        )

        try:
            # Step 1: Extract image bytes (placeholder - would need PDF library)
            # For now, we'll note where images would be extracted
            processed_images = []

            for img_info in enriched_doc.images:
                logger.debug(f"Processing image {img_info.image_id} on page {img_info.page_number}")

                try:
                    # Step 2: Extract image bytes from document
                    # Note: Actual extraction would require PyPDF2 or similar
                    # For now, this is a placeholder
                    image_bytes = self._extract_image_bytes(
                        original_bytes, img_info, enriched_doc.doc_type
                    )

                    if not image_bytes:
                        logger.warning(f"Could not extract bytes for {img_info.image_id}")
                        processed_images.append(img_info)
                        continue

                    # Step 3: Apply OCR if enabled
                    if self.config.ocr_enabled:
                        img_info.ocr_text = self._apply_ocr_to_image(image_bytes)

                    # Step 4: Generate description
                    if self.config.summary_method == "llm":
                        img_info.description = self._generate_image_description(
                            image_bytes, img_info.ocr_text
                        )
                    elif self.config.summary_method == "extractive":
                        # Use OCR text as description
                        img_info.description = img_info.ocr_text[: self.config.summary_max_length]

                    # Step 5: Classify image type
                    img_info.image_type = self._classify_image_type(
                        img_info.description, img_info.ocr_text
                    )

                    processed_images.append(img_info)
                    logger.debug(
                        f"Processed {img_info.image_id}: type={img_info.image_type}, "
                        f"description={img_info.description[:50]}..."
                    )

                except Exception as e:
                    logger.error(f"Error processing image {img_info.image_id}: {str(e)}")
                    processed_images.append(img_info)  # Keep original

            # Update document with processed images
            enriched_doc.images = processed_images

            # Step 6: Enrich markdown with image context
            if self.config.include_context_in_chunks:
                enriched_doc.markdown = self._enrich_markdown_with_images(
                    enriched_doc.markdown, processed_images
                )

            # Step 7: Update page-level image information
            self._update_page_images(enriched_doc)

            logger.info(
                f"✅ Processed {len(processed_images)} images for {enriched_doc.filename}"
            )

            return enriched_doc

        except Exception as e:
            logger.error(f"Failed to process images for {enriched_doc.filename}: {str(e)}")
            return enriched_doc  # Return original on error

    def _extract_image_bytes(
        self, document_bytes: bytes, image_info: ImageInfo, doc_type
    ) -> bytes | None:
        """
        Extract image bytes from document.

        Implementation approach:
        - IMAGE: return original bytes (standalone image files)
        - PDF: try PyMuPDF (fitz) if available; crop page region or extract largest image
        - DOCX: try python-docx if available; extract image blobs from relationships
        - PPTX: try python-pptx if available; extract image blobs from picture shapes

        Args:
            document_bytes: Original document bytes
            image_info: Information about image location
            doc_type: Document type

        Returns:
            bytes: Image bytes, or None if extraction fails
        """
        try:
            # Standalone image files: return bytes directly
            if doc_type == DocumentType.IMAGE:
                return document_bytes

            # PDF handling via PyMuPDF (fitz)
            if doc_type == DocumentType.PDF:
                try:
                    import fitz  # PyMuPDF

                    doc = fitz.open(stream=document_bytes, filetype="pdf")
                    page_index = max(0, (image_info.page_number or 1) - 1)
                    if page_index >= doc.page_count:
                        logger.warning(
                            f"PDF page index {page_index} out of range (count={doc.page_count})"
                        )
                        return None

                    page = doc.load_page(page_index)

                    # If position is known, crop that region from rendered page
                    if image_info.position:
                        pos = image_info.position
                        # Coordinates assumed in PDF user-space units; fall back gracefully
                        rect = fitz.Rect(pos.get("x", 0), pos.get("y", 0),
                                         pos.get("x", 0) + pos.get("width", 0),
                                         pos.get("y", 0) + pos.get("height", 0))
                        if rect.width > 0 and rect.height > 0:
                            pix = page.get_pixmap(clip=rect)
                            return pix.tobytes(output="png")

                    # Fallback: extract largest embedded image on the page
                    images = page.get_images(full=True)
                    if images:
                        # images: list of tuples (..., xref, ... width, height ...)
                        # Sort by area descending when size info available
                        try:
                            images_sorted = sorted(images, key=lambda i: (i[2] or 0) * (i[3] or 0), reverse=True)
                        except Exception:
                            images_sorted = images
                        for img in images_sorted:
                            xref = img[0] if isinstance(img[0], int) else img[1]
                            try:
                                extracted = doc.extract_image(xref)
                                if extracted and extracted.get("image"):
                                    return extracted["image"]
                            except Exception as e:
                                logger.debug(f"PDF image xref {xref} extract failed: {e}")

                    # Final fallback: render full page and return bytes (not ideal)
                    pix = page.get_pixmap()
                    return pix.tobytes(output="png")

                except ImportError:
                    logger.info("PyMuPDF (fitz) not installed; skipping PDF image extraction.")
                    return None
                except Exception as e:
                    logger.warning(f"PDF image extraction error: {e}")
                    return None

            # DOCX handling via python-docx
            if doc_type == DocumentType.DOCX:
                try:
                    from docx import Document  # python-docx
                    docx_obj = Document(BytesIO(document_bytes))
                    # Collect image blobs from relationships
                    image_blobs = []
                    for rel in docx_obj.part.rels.values():
                        reltype = getattr(rel, "reltype", "")
                        if reltype and reltype.endswith("/image"):
                            try:
                                image_blobs.append(rel.target_part.blob)
                            except Exception:
                                pass
                    if image_blobs:
                        # Heuristic: return first image (no reliable page mapping in DOCX)
                        return image_blobs[0]
                    logger.debug("DOCX contained no extractable image relationships.")
                    return None
                except ImportError:
                    logger.info("python-docx not installed; skipping DOCX image extraction.")
                    return None
                except Exception as e:
                    logger.warning(f"DOCX image extraction error: {e}")
                    return None

            # PPTX handling via python-pptx
            if doc_type == DocumentType.PPTX:
                try:
                    from pptx import Presentation  # python-pptx
                    prs = Presentation(BytesIO(document_bytes))
                    slide_index = max(0, (image_info.page_number or 1) - 1)
                    # Prefer image from specific slide when available
                    slides = list(prs.slides)
                    candidate_slides = [slides[slide_index]] if slide_index < len(slides) else slides
                    for slide in candidate_slides:
                        for shape in slide.shapes:
                            try:
                                # picture shapes expose .image.blob
                                if hasattr(shape, "image") and shape.image is not None:
                                    return shape.image.blob
                            except Exception:
                                continue
                    # Fallback: scan all slides
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            try:
                                if hasattr(shape, "image") and shape.image is not None:
                                    return shape.image.blob
                            except Exception:
                                continue
                    logger.debug("PPTX contained no extractable picture shapes.")
                    return None
                except ImportError:
                    logger.info("python-pptx not installed; skipping PPTX image extraction.")
                    return None
                except Exception as e:
                    logger.warning(f"PPTX image extraction error: {e}")
                    return None

            # Other types: not supported for embedded image extraction
            logger.debug(f"No image extraction strategy for doc_type={doc_type}")
            return None

        except Exception as e:
            logger.warning(f"Image extraction failure: {e}")
            return None

    def _apply_ocr_to_image(self, image_bytes: bytes) -> str:
        """
        Apply OCR to extract text from image.

        Args:
            image_bytes: Image content

        Returns:
            str: Extracted text
        """
        try:
            # Use Document Intelligence for OCR
            # Convert bytes to BytesIO
            image_stream = BytesIO(image_bytes)

            # Call Document Intelligence (Read model is good for OCR)
            poller = self.doc_intel_service.client.begin_analyze_document(
                "prebuilt-read", image_stream
            )
            result = poller.result()

            # Extract all text
            text_parts = []
            if hasattr(result, "content") and result.content:
                text_parts.append(result.content)

            ocr_text = " ".join(text_parts).strip()
            logger.debug(f"OCR extracted {len(ocr_text)} characters")

            return ocr_text

        except Exception as e:
            logger.warning(f"OCR failed: {str(e)}")
            return ""

    def _generate_image_description(self, image_bytes: bytes, ocr_text: str) -> str:
        """
        Generate description for image using GPT-4 Vision.

        Args:
            image_bytes: Image content
            ocr_text: OCR text from image

        Returns:
            str: Image description
        """
        try:
            description = self.openai_service.generate_image_description(
                image_bytes, ocr_text, max_length=self.config.summary_max_length
            )
            return description

        except Exception as e:
            logger.warning(f"Image description generation failed: {str(e)}")
            return ocr_text[: self.config.summary_max_length] if ocr_text else "Image"

    def _classify_image_type(self, description: str, ocr_text: str) -> str:
        """
        Classify image type from description and OCR text.

        Args:
            description: Image description
            ocr_text: OCR text

        Returns:
            str: Image type (chart, diagram, photo, table, etc.)
        """
        description_lower = description.lower()
        ocr_lower = ocr_text.lower()

        # Check for common image types
        if any(word in description_lower for word in ["chart", "graph", "plot"]):
            return "chart"
        elif any(word in description_lower for word in ["diagram", "flowchart", "flow chart"]):
            return "diagram"
        elif any(word in description_lower for word in ["table", "grid"]):
            return "table"
        elif any(word in description_lower for word in ["screenshot", "screen shot"]):
            return "screenshot"
        elif any(word in description_lower for word in ["photo", "photograph", "picture"]):
            return "photo"
        elif any(word in description_lower for word in ["logo", "icon"]):
            return "logo"
        elif any(word in description_lower for word in ["form", "application"]):
            return "form"
        else:
            return "image"

    def _enrich_markdown_with_images(self, markdown: str, images: list[ImageInfo]) -> str:
        """
        Enrich markdown content with image descriptions.

        Inserts image information at appropriate locations in the markdown.

        Args:
            markdown: Original markdown content
            images: List of processed images

        Returns:
            str: Enriched markdown with image context
        """
        if not images:
            return markdown

        enriched_parts = [markdown]

        # For each image, add a descriptive block
        # Note: Proper implementation would insert at exact locations based on page_number
        # For now, append image information at the end of relevant pages

        for img in images:
            if img.description or img.ocr_text:
                image_block = self._create_image_markdown_block(img)

                # In a full implementation, we'd insert this at the right location
                # For now, append to the end (will be improved with proper page tracking)
                enriched_parts.append(image_block)

        return "\n\n".join(enriched_parts)

    def _create_image_markdown_block(self, image_info: ImageInfo) -> str:
        """
        Create markdown block for an image.

        Args:
            image_info: Image information

        Returns:
            str: Markdown formatted image block
        """
        block_parts = []

        # Image reference
        block_parts.append(f"![Image: {image_info.image_type}]")
        block_parts.append("")

        # Description
        if image_info.description:
            block_parts.append(f"**Image Description**: {image_info.description}")

        # OCR text (if different from description and not too long)
        if image_info.ocr_text and image_info.ocr_text != image_info.description:
            ocr_preview = (
                image_info.ocr_text[:200] + "..."
                if len(image_info.ocr_text) > 200
                else image_info.ocr_text
            )
            block_parts.append(f"**Extracted Text**: {ocr_preview}")

        # Page reference
        if image_info.page_number:
            block_parts.append(f"*(Page {image_info.page_number})*")

        block_parts.append("")  # Empty line

        return "\n".join(block_parts)

    def _update_page_images(self, enriched_doc: EnrichedDocument):
        """
        Update page-level image information.

        Links images to their corresponding pages.

        Args:
            enriched_doc: Document to update
        """
        # Group images by page
        images_by_page: dict[int, list[ImageInfo]] = {}

        for img in enriched_doc.images:
            if img.page_number:
                if img.page_number not in images_by_page:
                    images_by_page[img.page_number] = []
                images_by_page[img.page_number].append(img)

        # Update pages
        for page in enriched_doc.pages:
            if page.page_number in images_by_page:
                page.images = images_by_page[page.page_number]


# Example usage
if __name__ == "__main__":
    import sys

    from src.config import get_config

    # Setup logging
    logging.basicConfig(
        level=logging.INFO, format="%(asctime)s - %(name)s - %(levelname)s - %(message)s"
    )

    try:
        config = get_config()

        # Initialize required services
        doc_intel_service = DocumentIntelligenceService(config.doc_intelligence)
        openai_service = OpenAIService(config.azure_openai, config.azure_openai_vision)
        image_service = ImageProcessingService(
            config.image_processing, doc_intel_service, openai_service
        )

        print("🖼️ Image Processing Service")
        print(f"Extract embedded images: {config.image_processing.extract_embedded_images}")
        print(f"OCR enabled: {config.image_processing.ocr_enabled}")
        print(f"Summary method: {config.image_processing.summary_method}")
        print(f"Include context in chunks: {config.image_processing.include_context_in_chunks}")

        print("\n✅ Image Processing Service initialized successfully!")
        print("\nNote: Full image extraction requires PDF/Office document libraries.")
        print("To be implemented: PyPDF2, python-docx, python-pptx")

    except Exception as e:
        print(f"Error: {str(e)}", file=sys.stderr)
        import traceback

        traceback.print_exc()
        sys.exit(1)
